"""
GFCC-style full Mass deck (matches parish PDF layout: dark slides, gold/white roles,
community footer, poster dividers). Readings filled from API/USCCB when available.

1920×1080 landscape. Accent color from liturgical calendar (replaces fixed gold where noted).
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any, List, Mapping, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from . import gfcc_flow_content as GFCC

SLIDE_WIDTH = Inches(20)
SLIDE_HEIGHT = Inches(11.25)

MARGIN_SIDE = Inches(0.75)
MARGIN_TOP = Inches(0.5)

_BG = RGBColor(18, 18, 22)
_GOLD_FALLBACK = RGBColor(220, 170, 90)
_BODY = RGBColor(245, 245, 245)
_MUTED = RGBColor(155, 155, 165)

_TITLE_PT = 38
_SECTION_PT = 30
_BODY_PT = 19
_META_PT = 14
_GREET_PT = 21
_FOOTER_PT = 13

_MAX_CHARS_READING = 900
_MAX_MARKED_BODY = 2600

_PROJECT_ROOT = Path(__file__).resolve().parents[1]
_OUTPUT_DIR = _PROJECT_ROOT / "outputs"

_COMMUNITY = "GWANGJU FILIPINO CATHOLIC COMMUNITY"


def _accent(liturgical_color: Optional[Mapping[str, Any]]) -> RGBColor:
    if liturgical_color and "rgb" in liturgical_color:
        r, g, b = liturgical_color["rgb"]
        return RGBColor(int(r), int(g), int(b))
    return _GOLD_FALLBACK


def _layout_blank(prs: Presentation):
    for layout in prs.slide_layouts:
        if "blank" in (layout.name or "").lower():
            return layout
    return prs.slide_layouts[-1]


def _set_slide_bg(slide, rgb: RGBColor):
    fi = slide.background.fill
    fi.solid()
    fi.fore_color.rgb = rgb


def _prep_tf(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)


def _style_para(p, *, size_pt, color, bold=False, italic=False, font_name="Calibri"):
    p.font.name = font_name
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color


def _add_community_footer(slide, footer_section: str, accent: RGBColor):
    lx = MARGIN_SIDE
    w = SLIDE_WIDTH - 2 * MARGIN_SIDE
    y = SLIDE_HEIGHT - Inches(0.95)
    foot = slide.shapes.add_textbox(lx, y, w, Inches(0.85))
    tf = foot.text_frame
    _prep_tf(tf)
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = _COMMUNITY
    _style_para(p0, size_pt=_FOOTER_PT, color=_MUTED, bold=True)
    p1 = tf.add_paragraph()
    p1.text = footer_section
    _style_para(p1, size_pt=_FOOTER_PT - 1, color=accent, bold=False)
    p1.space_before = Pt(2)


def _parse_marked_lines(marked: str) -> List[Tuple[str, str]]:
    out: List[Tuple[str, str]] = []
    for raw in (marked or "").split("\n"):
        raw = raw.strip()
        if not raw:
            continue
        role = "plain"
        if raw.startswith("<<P>>"):
            role, raw = "priest", raw[5:].strip()
        elif raw.startswith("<<A>>"):
            role, raw = "all", raw[5:].strip()
        elif raw.startswith("<<D>>"):
            role, raw = "direction", raw[5:].strip()
        elif raw.startswith("<<H>>"):
            role, raw = "hymn", raw[5:].strip()
        out.append((role, raw))
    return out


def _add_marked_slide(prs: Presentation, footer_section: str, marked_text: str, accent: RGBColor) -> None:
    slide = prs.slides.add_slide(_layout_blank(prs))
    _set_slide_bg(slide, _BG)
    lx, top, w = MARGIN_SIDE, MARGIN_TOP, SLIDE_WIDTH - 2 * MARGIN_SIDE
    body_h = SLIDE_HEIGHT - top - Inches(1.1)

    box = slide.shapes.add_textbox(lx, top, w, body_h)
    tf = box.text_frame
    _prep_tf(tf)
    tf.clear()
    first = True
    for role, line in _parse_marked_lines(marked_text):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        if role == "priest":
            _style_para(p, size_pt=_BODY_PT + 1, color=accent, bold=True)
        elif role == "all":
            _style_para(p, size_pt=_BODY_PT + 1, color=_BODY, bold=True)
        elif role == "direction":
            _style_para(p, size_pt=_META_PT + 1, color=accent, bold=False, italic=True)
        elif role == "hymn":
            _style_para(p, size_pt=_BODY_PT, color=_BODY, bold=True)
        else:
            _style_para(p, size_pt=_BODY_PT, color=_BODY, bold=False)
        p.space_after = Pt(5)

    _add_community_footer(slide, footer_section, accent)


def _chunk_marked_body(marked: str, limit: int = _MAX_MARKED_BODY) -> List[str]:
    if len(marked) <= limit:
        return [marked]
    parts, buf, n = [], [], 0
    for line in marked.split("\n"):
        line_len = len(line) + 1
        if n + line_len > limit and buf:
            parts.append("\n".join(buf))
            buf, n = [line], line_len
        else:
            buf.append(line)
            n += line_len
    if buf:
        parts.append("\n".join(buf))
    return parts if parts else [marked[:limit]]


def _add_marked_chunked(prs: Presentation, footer: str, marked: str, accent: RGBColor) -> None:
    chunks = _chunk_marked_body(marked)
    for i, ch in enumerate(chunks):
        foot = footer if len(chunks) == 1 else f"{footer} ({i + 1}/{len(chunks)})"
        _add_marked_slide(prs, foot, ch, accent)


def _add_divider_cover(
    prs: Presentation,
    *,
    celebrant: str,
    date: str,
    season: str,
    lectionary_cycle: str,
    gospel_reference: str,
    gospel_quote: str,
    quote_max_chars: int,
    accent: RGBColor,
) -> None:
    slide = prs.slides.add_slide(_layout_blank(prs))
    _set_slide_bg(slide, _BG)
    lx = MARGIN_SIDE
    lw = SLIDE_WIDTH - 2 * MARGIN_SIDE
    ty = MARGIN_TOP + Inches(0.5)

    g_line = (gospel_quote or "").strip()
    if quote_max_chars and len(g_line) > quote_max_chars:
        g_line = g_line[: quote_max_chars - 1].rstrip() + "\u2026"
    gref = (gospel_reference or "").strip() or "—"

    lines = [
        "MASS CELEBRANT:",
        celebrant,
        "",
        "\n".join(_COMMUNITY.split()),
        "",
        f"Gospel ({gref})",
    ]
    if g_line:
        lines.append(f"\u201c{g_line}\u201d")
    lines.extend(["", f"YEAR {(lectionary_cycle or '—').strip().upper()}", f"{date} · {(season or '').strip()}"])

    blk = slide.shapes.add_textbox(lx, ty, lw, Inches(6.5))
    tf = blk.text_frame
    _prep_tf(tf)
    tf.clear()
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        bold = "CELEBRANT" in line or line.startswith("MASS")
        _style_para(p, size_pt=_GREET_PT, color=_BODY if not bold else accent, bold=bold)
        p.space_after = Pt(3)

    _add_community_footer(slide, "Mass poster / divider", accent)


def _add_section_card(prs: Presentation, big_lines: str, footer_section: str, accent: RGBColor) -> None:
    slide = prs.slides.add_slide(_layout_blank(prs))
    _set_slide_bg(slide, _BG)
    lx, top, w = MARGIN_SIDE, MARGIN_TOP + Inches(1.15), SLIDE_WIDTH - 2 * MARGIN_SIDE
    box = slide.shapes.add_textbox(lx, top, w, Inches(4.5))
    tf = box.text_frame
    _prep_tf(tf)
    tf.clear()
    p = tf.paragraphs[0]
    p.text = big_lines
    _style_para(p, size_pt=44, color=accent, bold=True)
    p.alignment = PP_ALIGN.CENTER
    _add_community_footer(slide, footer_section, accent)


def chunk_plain_text(text: str, limit: int = _MAX_CHARS_READING) -> List[str]:
    if not (text or "").strip():
        return []
    norm = " ".join(text.split())
    if len(norm) <= limit:
        return [norm]
    sentences = re.split(r"(?<=[.!?])\s+", norm)
    out: List[str] = []
    buf = ""
    for s in sentences:
        w = s.strip()
        if not w:
            continue
        spacer = " " if buf else ""
        if len(buf) + len(spacer) + len(w) <= limit:
            buf += spacer + w
        else:
            if buf:
                out.append(buf.strip())
            if len(w) <= limit:
                buf = w
            else:
                for i in range(0, len(w), limit):
                    piece = w[i : i + limit].strip()
                    if piece:
                        out.append(piece)
                buf = ""
    if buf:
        out.append(buf.strip())
    return out if out else [norm[:limit]]


def _paragraphs(tf, *, size_pt, color, bold=False):
    tf.clear()
    p = tf.paragraphs[0]
    _style_para(p, size_pt=size_pt, color=color, bold=bold)


def _fill_multipara(tf, text: str, *, size_pt: int, color: RGBColor):
    tf.clear()
    raw = (text or "").strip()
    parts = [b.strip() for b in raw.split("\n\n") if b.strip()] or ([raw] if raw else [""])
    first = True
    for block in parts:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = block
        _style_para(p, size_pt=size_pt, color=color)
        p.space_after = Pt(5)


def _add_reading_block(
    prs: Presentation,
    *,
    section: str,
    reference: str,
    body: str,
    unavailable_note: str,
    lotw_banner: bool,
    footer_tag: str,
    accent: RGBColor,
) -> None:
    ref = (reference or "").strip() or "—"
    body = (body or "").strip()

    def one_slide(head: str, sub: str, main: str):
        slide = prs.slides.add_slide(_layout_blank(prs))
        _set_slide_bg(slide, _BG)
        lx, top, w = MARGIN_SIDE, MARGIN_TOP, SLIDE_WIDTH - 2 * MARGIN_SIDE
        title_h = Inches(1.12) if lotw_banner else Inches(0.92)
        title_box = slide.shapes.add_textbox(lx, top, w, title_h)
        tf_t = title_box.text_frame
        _prep_tf(tf_t)
        tf_t.clear()
        if lotw_banner:
            p0 = tf_t.paragraphs[0]
            p0.text = "Liturgy of the Word"
            _style_para(p0, size_pt=_SECTION_PT - 2, color=accent, bold=True)
            p1 = tf_t.add_paragraph()
            p1.text = head if "continued" in head.lower() else f"{section} ({ref})"
            _style_para(p1, size_pt=_META_PT + 2, color=_MUTED, bold=False)
        else:
            _paragraphs(tf_t, size_pt=_SECTION_PT, color=accent, bold=True)
            tf_t.paragraphs[0].text = head

        sub_top = top + title_h + Inches(0.06)
        sub_h = Inches(0.48)
        sub_box = slide.shapes.add_textbox(lx, sub_top, w, sub_h)
        _prep_tf(sub_box.text_frame)
        _paragraphs(sub_box.text_frame, size_pt=_META_PT, color=_MUTED)
        sub_box.text_frame.paragraphs[0].text = sub

        body_top = sub_top + sub_h + Inches(0.12)
        body_h = SLIDE_HEIGHT - body_top - Inches(1.0)
        bsh = slide.shapes.add_textbox(lx, body_top, w, body_h)
        _prep_tf(bsh.text_frame)
        _paragraphs(bsh.text_frame, size_pt=_BODY_PT, color=_BODY)
        _fill_multipara(bsh.text_frame, main, size_pt=_BODY_PT, color=_BODY)
        _add_community_footer(slide, footer_tag, accent)

    if not body:
        one_slide(section, ref, unavailable_note)
        return
    chunks = chunk_plain_text(body)
    total = len(chunks)
    for i, chunk in enumerate(chunks):
        head = section if i == 0 else f"{section} (continued)"
        if lotw_banner:
            sub = "" if total <= 1 else f"Slide {i + 1} of {total}"
        else:
            sub = ref if total <= 1 else f"{ref}  ·  slide {i + 1} of {total}"
        one_slide(head, sub, chunk)


def _add_title_slide(
    prs: Presentation,
    *,
    title: str,
    date: str,
    celebrant: str,
    gospel_reference: str,
    gospel_quote: str,
    season: str,
    lectionary_cycle: str,
    liturgical_color: Optional[Mapping[str, Any]],
    quote_attribution: Optional[str],
    quote_max_chars: int,
    accent: RGBColor,
) -> None:
    slide = prs.slides.add_slide(_layout_blank(prs))
    _set_slide_bg(slide, _BG)
    lx, w = MARGIN_SIDE, SLIDE_WIDTH - 2 * MARGIN_SIDE
    y = MARGIN_TOP + Inches(0.35)

    tb = slide.shapes.add_textbox(lx, y, w, Inches(1.05))
    tft = tb.text_frame
    _prep_tf(tft)
    tft.clear()
    p0 = tft.paragraphs[0]
    p0.text = title or "Mass"
    _style_para(p0, size_pt=_TITLE_PT, color=accent, bold=True)

    g_line = (gospel_quote or "").strip()
    if quote_max_chars and len(g_line) > quote_max_chars:
        g_line = g_line[: quote_max_chars - 1].rstrip() + "\u2026"
    gref = (gospel_reference or "").strip() or "—"

    meta = (
        f"{date}\n\nCelebrant: {celebrant}\n\n"
        f"Gospel: {gref}\n"
        f"Season: {(season or '—').strip()} · Sunday Lectionary Year {(lectionary_cycle or '—').strip().upper()}"
    )
    if g_line:
        meta += f"\n\nExcerpt:\n\u201c{g_line}\u201d"
    if liturgical_color:
        meta += f"\n\nLiturgical color: {liturgical_color.get('color_name', '')} ({liturgical_color.get('season', '')})"

    mb = slide.shapes.add_textbox(lx, y + Inches(1.15), w, Inches(3.4))
    _prep_tf(mb.text_frame)
    _fill_multipara(mb.text_frame, meta, size_pt=_GREET_PT, color=_BODY)

    if quote_attribution and g_line:
        nb = slide.shapes.add_textbox(lx, SLIDE_HEIGHT - Inches(1.2), w, Inches(0.75))
        _prep_tf(nb.text_frame)
        _fill_multipara(nb.text_frame, str(quote_attribution), size_pt=_META_PT, color=_MUTED)

    _add_community_footer(slide, "Title", accent)


def generate_mass_ppt(
    title: str,
    gospel_reference: str,
    gospel_quote: str,
    season: str,
    lectionary_cycle: str,
    celebrant: str,
    date: str,
    *,
    gospel_full_text: str = "",
    first_reading_ref: str = "",
    first_reading_text: str = "",
    psalm_ref: str = "",
    psalm_text: str = "",
    second_reading_ref: str = "",
    second_reading_text: str = "",
    quote_attribution=None,
    quote_max_chars: int = 400,
    liturgical_color: Optional[Mapping[str, Any]] = None,
):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    accent = _accent(liturgical_color)

    g_line = (gospel_quote or "").strip()
    if quote_max_chars and len(g_line) > quote_max_chars:
        g_line = g_line[: quote_max_chars - 1].rstrip() + "\u2026"

    unavail = (
        "Full text was not loaded from bible.usccb.org. "
        "Open today’s readings for this date and paste if needed."
    )

    ctx = dict(
        celebrant=celebrant,
        date=date,
        season=season,
        lectionary_cycle=lectionary_cycle,
        gospel_reference=gospel_reference or "",
        gospel_quote=g_line,
        quote_max_chars=quote_max_chars,
        accent=accent,
    )

    # --- Pre-Mass (GFCC PDF p.1 style) ---
    _add_marked_slide(prs, "Pre-Mass", GFCC.SILENT_REMINDER, accent)

    _add_title_slide(
        prs,
        title=title,
        date=date,
        celebrant=celebrant,
        gospel_reference=gospel_reference or "",
        gospel_quote=g_line,
        season=season,
        lectionary_cycle=lectionary_cycle,
        liturgical_color=liturgical_color,
        quote_attribution=quote_attribution,
        quote_max_chars=quote_max_chars,
        accent=accent,
    )

    _add_marked_chunked(prs, "Entrance", GFCC.ENTRANCE_HYMN_1 + "\n" + GFCC.ENTRANCE_HYMN_2, accent)
    _add_divider_cover(prs, **ctx)

    # --- Introductory Rites ---
    _add_marked_slide(prs, "Introductory Rites", GFCC.SIGN_CROSS, accent)
    _add_marked_slide(prs, "Penitential Act", GFCC.GREETING_EXTENDED, accent)
    _add_marked_slide(prs, "Penitential Act", GFCC.CONFITEOR_OPEN, accent)
    _add_marked_slide(prs, "Penitential Act", GFCC.ABSOLUTION_PENITENTIAL, accent)
    _add_marked_slide(prs, "Kyrie Eleison", GFCC.KYRIE, accent)
    _add_marked_chunked(prs, "Gloria", GFCC.GLORIA_1 + "\n" + GFCC.GLORIA_2 + "\n" + GFCC.GLORIA_3, accent)
    _add_marked_slide(prs, "Liturgy of the Word", GFCC.OPENING_PRAYER, accent)

    # --- Liturgy of the Word ---
    _add_section_card(prs, "LITURGY OF\nTHE WORD", "Liturgy of the Word", accent)

    _add_reading_block(
        prs,
        section="First Reading",
        reference=first_reading_ref or "—",
        body=(first_reading_text or "").strip(),
        unavailable_note=unavail,
        lotw_banner=True,
        footer_tag="Liturgy of the Word",
        accent=accent,
    )
    _add_reading_block(
        prs,
        section="Responsorial Psalm",
        reference=psalm_ref or "—",
        body=(psalm_text or "").strip(),
        unavailable_note=unavail,
        lotw_banner=True,
        footer_tag="Liturgy of the Word",
        accent=accent,
    )
    if (second_reading_ref or "").strip():
        _add_reading_block(
            prs,
            section="Second Reading",
            reference=second_reading_ref.strip(),
            body=(second_reading_text or "").strip(),
            unavailable_note=unavail,
            lotw_banner=True,
            footer_tag="Liturgy of the Word",
            accent=accent,
        )

    _add_marked_slide(prs, "Gospel Acclamation", GFCC.ALLELUIA_SING, accent)
    _add_marked_slide(prs, "Gospel Acclamation", GFCC.ALLELUIA_COMMENTATOR, accent)
    _add_marked_slide(prs, "Gospel Acclamation", GFCC.GOSPEL_INTRO, accent)

    g_body = ((gospel_full_text or "").strip() or (gospel_quote or "").strip())
    _add_reading_block(
        prs,
        section="Gospel",
        reference=gospel_reference or "—",
        body=g_body,
        unavailable_note=unavail,
        lotw_banner=False,
        footer_tag="Gospel",
        accent=accent,
    )

    _add_marked_slide(prs, "Gospel Acclamation", GFCC.GOSPEL_END, accent)
    _add_marked_slide(
        prs,
        "Homily",
        "<<D>>Time for the homily — Father will now preach.\n<<D>>Commentator may introduce the theme.",
        accent,
    )
    _add_divider_cover(prs, **ctx)

    # --- Creed ---
    _add_marked_chunked(prs, "Nicene Creed", GFCC.CREED_1 + "\n" + GFCC.CREED_2 + "\n" + GFCC.CREED_3, accent)
    _add_divider_cover(prs, **ctx)

    # --- Prayer of the Faithful ---
    _add_marked_slide(prs, "Prayer of the Faithful", GFCC.PRAYER_FAITHFUL_1, accent)
    _add_marked_slide(prs, "Prayer of the Faithful", GFCC.PRAYER_FAITHFUL_2, accent)
    _add_divider_cover(prs, **ctx)

    # --- Liturgy of the Eucharist ---
    _add_marked_chunked(prs, "Liturgy of the Eucharist", GFCC.OFFERTORY_HYMN, accent)
    _add_section_card(prs, "LITURGY OF\nTHE EUCHARIST", "Liturgy of the Eucharist", accent)
    _add_marked_slide(prs, "Liturgy of the Eucharist", GFCC.PRAY_BRETHREN, accent)
    _add_section_card(prs, "LITURGY OF\nTHE EUCHARIST", "Liturgy of the Eucharist", accent)
    _add_marked_slide(prs, "Liturgy of the Eucharist", GFCC.PREFACE_DIALOGUE, accent)
    _add_marked_slide(prs, "Liturgy of the Eucharist", GFCC.PREFACE_ACCLAIM, accent)
    _add_marked_chunked(prs, "Sanctus", GFCC.SANCTUS, accent)
    _add_section_card(prs, "LITURGY OF\nTHE EUCHARIST", "Liturgy of the Eucharist", accent)
    _add_marked_slide(prs, "The Eucharistic Prayer", GFCC.MYSTERY_FAITH, accent)
    _add_section_card(prs, "LITURGY OF\nTHE EUCHARIST", "Liturgy of the Eucharist", accent)
    _add_marked_slide(prs, "Great Amen", GFCC.GREAT_AMEN, accent)
    _add_marked_slide(prs, "Our Father", GFCC.OUR_FATHER_KO_1, accent)
    _add_marked_slide(prs, "Our Father", GFCC.OUR_FATHER_KO_2, accent)
    _add_divider_cover(prs, **ctx)
    _add_marked_slide(prs, "The Communion Rite", GFCC.COMMUNION_RITE_DELIVER, accent)
    _add_divider_cover(prs, **ctx)
    _add_marked_slide(prs, "Sign of Peace", GFCC.SIGN_PEACE, accent)
    _add_marked_slide(prs, "Lamb of God", GFCC.LAMB_OF_GOD, accent)
    _add_marked_slide(prs, "The Communion Rite", GFCC.COMMUNION_DIALOGUE, accent)
    _add_divider_cover(prs, **ctx)
    _add_marked_chunked(prs, "Communion", GFCC.COMMUNION_HYMN, accent)
    _add_marked_slide(prs, "The Communion Rite", GFCC.POST_COMMUNION, accent)
    _add_divider_cover(prs, **ctx)

    # --- Announcements ---
    _add_marked_slide(prs, "Announcements", GFCC.ANNOUNCEMENTS_TITLE, accent)
    _add_marked_slide(prs, "Announcements", GFCC.WELCOME_NEWCOMERS, accent)
    _add_marked_slide(prs, "Announcements", GFCC.CONFESSION_SLIDE, accent)
    _add_marked_slide(prs, "Announcements", GFCC.COLLECTION_PLACEHOLDER, accent)
    _add_marked_slide(prs, "Announcements", GFCC.SPONSORSHIP, accent)
    _add_marked_slide(prs, "Announcements", GFCC.FB_UPDATES, accent)

    _add_marked_slide(prs, "Final Blessing", GFCC.FINAL_BLESSING, accent)
    _add_marked_chunked(prs, "Recessional", GFCC.RECESSIONAL_1 + "\n" + GFCC.RECESSIONAL_2, accent)
    _add_divider_cover(prs, **ctx)

    if quote_attribution and g_line:
        _add_marked_slide(prs, "Scripture note", f"<<D>>{quote_attribution}", accent)

    _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    out = _OUTPUT_DIR / "mass_presentation.pptx"
    prs.save(out)
    print(f"✅ PowerPoint created: {out}")
