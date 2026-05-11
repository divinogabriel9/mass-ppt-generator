"""
Church Media Generator — minimal web UI + JSON API.

Run from project root:
  cd church_media_generator && uvicorn server:app --reload --host 127.0.0.1 --port 8000
"""

from __future__ import annotations

import io
import shutil
import subprocess
import zipfile
from pathlib import Path
from typing import Any, Optional

from fastapi import FastAPI, File, HTTPException, Request, UploadFile
from fastapi.responses import HTMLResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from pydantic import BaseModel, Field

from pipeline import (
    GenerationResult,
    PreviewPayload,
    fetch_preview,
    generate_mass_media,
    refresh_all_song_sections,
    refresh_song_section,
)
from services.community_config import (
    LOGO_RELATIVE,
    get_community_name,
    logo_file_absolute,
    update_community,
    uploads_dir,
)
from services.lyrics_fetcher import fetch_and_store_for_selection
from services.ppt_preview_render import render_ppt_preview_pngs
from services.song_catalog import import_song_rows, import_titles

# Optional outputs produced alongside mass_poster.png (Phase 3)
_BUNDLE_OPTIONAL = (
    "gospel_moment.png",
    "mass_poster_16x9.png",
    "mass_poster_instagram_square.png",
    "mass_poster_instagram_story.png",
    "mass_poster_open_graph.png",
)

_PROJECT = Path(__file__).resolve().parent
_OUTPUT_DIR = _PROJECT / "outputs"
_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
_PREVIEW_DIR = _OUTPUT_DIR / "preview_slides"
_PREVIEW_DIR.mkdir(parents=True, exist_ok=True)

_UPLOAD_DIR = uploads_dir()
_UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

_ALLOWED_LOGO_TYPES = frozenset(
    ("image/png", "image/jpeg", "image/webp", "image/gif", "image/x-png", "image/jpg")
)
_MAX_LOGO_BYTES = 2_500_000

_BUNDLE_NAME = "mass_bundle.zip"


def _ai_poster_download_urls() -> dict[str, str]:
    """Stable URLs under ``/media/posters/`` for Hugging Face layout exports."""
    base = _OUTPUT_DIR / "posters"
    out: dict[str, str] = {}
    mapping = {
        "instagram": "mass_poster_instagram.png",
        "story": "mass_poster_story.png",
        "facebook": "mass_poster_facebook.png",
    }
    for key, fname in mapping.items():
        p = base / fname
        if p.is_file():
            out[key] = f"/media/posters/{fname}"
    return out


def _latest_pptx_path() -> Optional[Path]:
    """Most recently modified deck in ``outputs/`` (supports stem-based filenames)."""
    cands = [p for p in _OUTPUT_DIR.glob("*.pptx") if p.is_file()]
    if not cands:
        return None
    return max(cands, key=lambda p: p.stat().st_mtime)


def _write_mass_bundle_zip(result: GenerationResult) -> None:
    """Pack generated PPT, posters, stem-based social PNGs, gospel art, and optional extras."""
    out = _OUTPUT_DIR / _BUNDLE_NAME
    entries: list[tuple[Path, str]] = []
    for attr in ("pptx_path", "poster_path", "poster_ppt_path"):
        p = getattr(result, attr, None)
        if p and Path(p).is_file():
            pp = Path(p)
            entries.append((pp, pp.name))
    if result.poster_path:
        parent = Path(result.poster_path).parent
        stem = Path(result.poster_path).stem
        for child in sorted(parent.glob(f"{stem}_*.png")):
            if child.is_file() and all(o[0] != child for o in entries):
                entries.append((child, child.name))
    if result.export_stem:
        g = _OUTPUT_DIR / f"{result.export_stem}_gospel_moment.png"
        if g.is_file() and all(o[0] != g for o in entries):
            entries.append((g, g.name))
    post_dir = _OUTPUT_DIR / "posters"
    if post_dir.is_dir():
        for child in sorted(post_dir.glob("*.png")):
            if child.is_file() and all(o[0] != child for o in entries):
                entries.append((child, f"posters/{child.name}"))
    for name in _BUNDLE_OPTIONAL:
        p = _OUTPUT_DIR / name
        if p.is_file() and all(o[0] != p for o in entries):
            entries.append((p, name))
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zf:
        for path, arc in entries:
            zf.write(path, arcname=arc)


app = FastAPI(title="Church Media Generator")
templates = Jinja2Templates(directory=str(_PROJECT / "templates"))
app.mount("/media", StaticFiles(directory=str(_OUTPUT_DIR)), name="media")
app.mount("/uploads", StaticFiles(directory=str(_UPLOAD_DIR)), name="uploads")
app.mount("/preview", StaticFiles(directory=str(_PREVIEW_DIR)), name="preview")


def _preview_to_json(p: PreviewPayload) -> dict[str, Any]:
    lc = p.liturgical_color
    liturgical: Optional[dict[str, Any]] = None
    if lc:
        liturgical = {
            "color_name": lc.get("color_name"),
            "hex": lc.get("hex"),
            "season": lc.get("season"),
            "rgb": list(lc.get("rgb", ())),
        }
    return {
        "ok": p.ok,
        "error": p.error,
        "title": p.title,
        "gospel_reference": p.gospel_reference,
        "season": p.season,
        "lectionary_cycle": p.lectionary_cycle,
        "liturgical_color": liturgical,
        "gospel_text_length": p.gospel_text_length,
        "sentences": p.sentences,
        "sentence_count": len(p.sentences),
        "quote_attribution": p.quote_attribution,
        "songs_by_section": p.songs_by_section,
        "gospel_quote": p.gospel_quote,
        "default_song_selections": p.default_song_selections,
        "estimated_slide_count": p.estimated_slide_count,
    }


class SongSelection(BaseModel):
    entrance: Optional[str] = None
    offertory: Optional[str] = None
    communion_1: Optional[str] = None
    communion_2: Optional[str] = None
    recessional: Optional[str] = None
    meditation: Optional[str] = None


class CommunityNameBody(BaseModel):
    community_name: str = Field(..., min_length=1, max_length=240)


class PreviewBody(BaseModel):
    date: str = Field(..., min_length=8, description="YYYY-MM-DD")


class GenerateBody(BaseModel):
    date: str = Field(..., min_length=8)
    celebrant: str = Field(..., min_length=1, max_length=200)
    sentence_index: Optional[int] = Field(None, ge=0)
    poster_template: str = Field(
        "liturgical_color",
        description="liturgical_color | classic_white",
    )
    include_social_exports: bool = Field(True)
    include_gospel_art: bool = Field(True)
    include_ai_mass_poster: bool = Field(
        True,
        description="Run Hugging Face AI poster export to outputs/posters/ (uses token if set).",
    )
    ai_poster_style: str = Field(
        "cinematic",
        max_length=64,
        description="HF hero art style key from data/styles.json (e.g. renaissance, stained_glass).",
    )
    community_name: Optional[str] = Field(None, max_length=240)
    songs: Optional[SongSelection] = None


class RefreshSongsBody(BaseModel):
    date: str = Field(..., min_length=8, description="YYYY-MM-DD")
    section: str = Field(..., min_length=3, max_length=40)
    current_ids: list[str] = Field(default_factory=list)


class RefreshAllSongsBody(BaseModel):
    date: str = Field(..., min_length=8, description="YYYY-MM-DD")
    current_ids: dict[str, list[str]] = Field(default_factory=dict)


class ImportSongsBody(BaseModel):
    entrance: list[str] = Field(default_factory=list)
    offertory: list[str] = Field(default_factory=list)
    communion: list[str] = Field(default_factory=list)
    recessional: list[str] = Field(default_factory=list)


class SongRowBody(BaseModel):
    title: str = Field(..., min_length=1, max_length=240)
    language: str = Field("English", max_length=40)
    mass_part: list[str] = Field(default_factory=list)


class ImportSongRowsBody(BaseModel):
    songs: list[SongRowBody] = Field(default_factory=list)


class LyricsSelectionBody(BaseModel):
    section: str = Field(..., min_length=3, max_length=40)
    id: str = Field(..., min_length=1, max_length=160)


class FetchLyricsBody(BaseModel):
    selections: list[LyricsSelectionBody] = Field(default_factory=list)


def _resolve_soffice_bin() -> Optional[str]:
    custom = shutil.which("soffice")
    if custom:
        return custom
    mac_bin = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if Path(mac_bin).is_file():
        return mac_bin
    return None


def _extract_ppt_text_slides(ppt: Path) -> list[dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError:
        return []
    prs = Presentation(str(ppt))
    slides: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        lines: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                txt = (shape.text or "").strip()
                if txt:
                    lines.append(txt)
        slides.append({"index": idx + 1, "text": "\n\n".join(lines)[:2000]})
    return slides


@app.post("/api/ppt-preview/refresh")
def api_ppt_preview_refresh() -> dict[str, Any]:
    """Render PPT slides to PNG images for in-app visual preview."""
    ppt = _latest_pptx_path()
    if not ppt or not ppt.is_file():
        return {"ok": True, "mode": "text", "slides": [], "message": "Generate deck first."}
    soffice = _resolve_soffice_bin()
    if not soffice:
        return {
            "ok": True,
            "mode": "text",
            "slides": _extract_ppt_text_slides(ppt),
            "message": "Install LibreOffice for exact image preview. Showing text fallback.",
        }

    for p in _PREVIEW_DIR.glob("*"):
        if p.is_file():
            p.unlink(missing_ok=True)

    png_paths, pdf_msg = render_ppt_preview_pngs(ppt, _PREVIEW_DIR, soffice_bin=soffice)
    if not png_paths:
        return {
            "ok": True,
            "mode": "text",
            "slides": _extract_ppt_text_slides(ppt),
            "message": (pdf_msg or "Could not render slide images.") + " Showing text fallback.",
        }

    ts = int(png_paths[0].stat().st_mtime) if png_paths else 0
    slides = [
        {
            "index": i + 1,
            "image_url": f"/preview/{f.name}?t={ts}&v={i}",
        }
        for i, f in enumerate(png_paths)
    ]
    msg = pdf_msg or ""
    if not msg.strip():
        msg = "Full-deck preview (PDF rasterization)."
    return {"ok": True, "mode": "image", "slides": slides, "message": msg}


@app.get("/health")
def health() -> dict[str, str]:
    return {"status": "ok"}


@app.get("/api/community")
def api_community() -> dict[str, Any]:
    logo_exists = logo_file_absolute().is_file()
    return {
        "community_name": get_community_name(),
        "logo_url": "/uploads/community_logo.png" if logo_exists else None,
    }


@app.post("/api/community")
def api_set_community_name(body: CommunityNameBody) -> dict[str, Any]:
    update_community(community_name=body.community_name.strip())
    return {"ok": True, "community_name": get_community_name()}


@app.post("/api/upload-logo")
async def api_upload_logo(file: UploadFile = File(...)) -> dict[str, Any]:
    ctype = (file.content_type or "").split(";")[0].strip().lower()
    if ctype not in _ALLOWED_LOGO_TYPES:
        raise HTTPException(
            status_code=400,
            detail="Use a PNG, JPEG, WebP, or GIF image.",
        )
    raw = await file.read()
    if len(raw) > _MAX_LOGO_BYTES:
        raise HTTPException(status_code=400, detail="Image must be at most about 2.5 MB.")

    try:
        from PIL import Image
    except ImportError as exc:
        raise HTTPException(status_code=500, detail="Pillow is required for image upload.") from exc

    try:
        im = Image.open(io.BytesIO(raw))
        im.load()
    except OSError as exc:
        raise HTTPException(status_code=400, detail="Could not read image file.") from exc

    if im.mode in ("RGBA", "LA") or (im.mode == "P" and "transparency" in getattr(im, "info", {})):
        im = im.convert("RGBA")
    else:
        im = im.convert("RGB")

    _UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    out_path = logo_file_absolute()
    im.save(out_path, format="PNG", optimize=True)

    update_community(logo_path=LOGO_RELATIVE)
    return {
        "ok": True,
        "logo_url": "/uploads/community_logo.png",
        "message": "Logo saved. It will appear on the next generated poster.",
    }


@app.get("/", response_class=HTMLResponse)
def index(request: Request) -> Any:
    # Starlette 0.28+: (request, name, context); request is injected into the template context.
    return templates.TemplateResponse(
        request,
        "index.html",
        {"title": "Church Media Generator"},
    )


@app.post("/api/preview")
def api_preview(body: PreviewBody) -> Any:
    p = fetch_preview(body.date.strip())
    return _preview_to_json(p)


@app.post("/api/generate")
def api_generate(body: GenerateBody) -> Any:
    song_map = body.songs.model_dump(exclude_none=True) if body.songs else None
    result = generate_mass_media(
        body.date.strip(),
        body.celebrant.strip(),
        sentence_index=body.sentence_index,
        poster_template=body.poster_template,
        include_social_exports=body.include_social_exports,
        include_gospel_art=body.include_gospel_art,
        include_ai_mass_poster=body.include_ai_mass_poster,
        ai_poster_style=body.ai_poster_style.strip() or "cinematic",
        community_name=body.community_name.strip() if body.community_name else None,
        song_selections=song_map,
    )
    if not result.ok:
        raise HTTPException(status_code=400, detail=result.error or "Generation failed.")
    _write_mass_bundle_zip(result)
    zip_ready = (_OUTPUT_DIR / _BUNDLE_NAME).is_file()
    lc = result.liturgical_color
    liturgical_payload: Optional[dict[str, Any]] = None
    if lc:
        liturgical_payload = {
            "color_name": lc.get("color_name"),
            "hex": lc.get("hex"),
            "season": lc.get("season"),
            "rgb": list(lc.get("rgb", ())),
        }
    poster_ppt = result.poster_ppt_path
    ppt_url = f"/media/{result.pptx_path.name}" if result.pptx_path and result.pptx_path.is_file() else None
    poster_url = f"/media/{result.poster_path.name}" if result.poster_path and result.poster_path.is_file() else None
    poster_ppt_url = (
        f"/media/{poster_ppt.name}" if poster_ppt and Path(poster_ppt).is_file() else None
    )
    payload: dict[str, Any] = {
        "ok": True,
        "title": result.title,
        "gospel_reference": result.gospel_reference,
        "slide_excerpt": result.slide_line_preview,
        "gospel_quote": result.gospel_quote,
        "liturgical_color": liturgical_payload,
        "selected_songs": result.selected_songs,
        "slide_count": result.slide_count,
        "export_stem": result.export_stem,
    }
    if ppt_url:
        payload["pptx_url"] = ppt_url
    if poster_url:
        payload["poster_url"] = poster_url
    if poster_ppt_url:
        payload["poster_ppt_url"] = poster_ppt_url
    payload["ai_poster_urls"] = _ai_poster_download_urls()
    return {
        **payload,
        **(
            {"zip_url": f"/media/{_BUNDLE_NAME}"}
            if zip_ready
            else {}
        ),
    }


@app.post("/api/songs/refresh")
def api_refresh_songs(body: RefreshSongsBody) -> dict[str, Any]:
    sec = body.section.strip().lower()
    if sec not in {"entrance", "offertory", "communion", "recessional", "meditation"}:
        raise HTTPException(status_code=400, detail="Invalid section.")
    songs = refresh_song_section(
        date=body.date.strip(),
        section=sec,
        current_ids=[str(x) for x in (body.current_ids or [])],
        limit=10,
    )
    return {"ok": True, "section": sec, "songs": songs}


@app.post("/api/songs/refresh-all")
def api_refresh_all_songs(body: RefreshAllSongsBody) -> dict[str, Any]:
    songs = refresh_all_song_sections(
        date=body.date.strip(),
        current_ids=body.current_ids or {},
        limit=10,
    )
    return {"ok": True, "songs_by_section": songs}


@app.post("/api/songs/import")
def api_import_songs(body: ImportSongsBody) -> dict[str, Any]:
    summary = import_titles(
        {
            "entrance": body.entrance,
            "offertory": body.offertory,
            "communion": body.communion,
            "recessional": body.recessional,
        }
    )
    return {"ok": True, **summary}


@app.post("/api/songs/import-list")
def api_import_song_rows(body: ImportSongRowsBody) -> dict[str, Any]:
    summary = import_song_rows([s.model_dump() for s in body.songs])
    return {"ok": True, **summary}


@app.post("/api/songs/fetch-lyrics")
def api_fetch_lyrics(body: FetchLyricsBody) -> dict[str, Any]:
    if not body.selections:
        return {"ok": True, "updated": 0, "skipped": 0, "results": []}
    results = [fetch_and_store_for_selection({"section": s.section, "id": s.id}) for s in body.selections]
    updated = sum(1 for r in results if r.get("ok") and r.get("reason") == "fetched")
    skipped = len(results) - updated
    return {"ok": True, "updated": updated, "skipped": skipped, "results": results}
